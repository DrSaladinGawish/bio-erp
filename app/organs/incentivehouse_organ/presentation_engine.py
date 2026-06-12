"""
P1-B3: Smart Presentation Engine
Auto-generates PowerPoint from event data — Zero Gap Compliance
"""

import os
from datetime import datetime
from typing import List, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from app.organs.incentivehouse_organ.models_production import (
    Event,
    Client,
    SalesLineItem,
)
from app.organs.incentivehouse_organ.db import get_sync_session_factory
from app.organs.incentivehouse_organ.rbac import Permission, require_permission


class PresentationEngine:
    PRIMARY = RGBColor(0x1A, 0x3A, 0x5C)
    ACCENT = RGBColor(0xE8, 0xB9, 0x23)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x0D, 0x21, 0x37)
    GRAY = RGBColor(0x6C, 0x75, 0x7D)

    def __init__(self, output_dir: str = "D:/ERP System/BIO_ERP/presentations"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate_event_deck(self, event_id: int, db_session) -> str:
        event = db_session.query(Event).get(event_id)
        if not event:
            raise ValueError(f"Event {event_id} not found")

        client = (
            db_session.query(Client).get(event.client_id) if event.client_id else None
        )
        line_items = (
            db_session.query(SalesLineItem)
            .filter(SalesLineItem.invoice_id == event_id)
            .all()
        )

        self._add_title_slide(event, client)
        self._add_event_overview(event, client)
        self._add_financial_summary(event, line_items)
        self._add_timeline_slide(event)
        self._add_vendor_breakdown(line_items)
        self._add_checklist_slide(event)
        self._add_closing_slide(event, client)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"IH_Event_{event_id}_{timestamp}.pptx"
        filepath = self.output_dir / filename
        self.prs.save(str(filepath))
        return str(filepath)

    def _add_title_slide(self, event: Event, client: Optional[Client]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.PRIMARY
        bg.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = event.event_name or "Event Presentation"
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        if client:
            client_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12), Inches(0.8)
            )
            ctf = client_box.text_frame
            ctf.text = f"Prepared for: {client.name}"
            cp = ctf.paragraphs[0]
            cp.font.size = Pt(24)
            cp.font.color.rgb = self.ACCENT

        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12), Inches(0.5)
        )
        dtf = date_box.text_frame
        dtf.text = datetime.now().strftime("%B %d, %Y")
        dp = dtf.paragraphs[0]
        dp.font.size = Pt(14)
        dp.font.color.rgb = self.GRAY

    def _add_event_overview(self, event: Event, client: Optional[Client]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY
        header.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.8)
        )
        tf = title.text_frame
        tf.text = "Event Overview"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        rows = [
            ["Event Name", event.event_name or "N/A"],
            ["Client", client.name if client else "N/A"],
            [
                "Event Date",
                event.event_date.strftime("%Y-%m-%d") if event.event_date else "TBD",
            ],
            ["Status", event.lifecycle_status or "draft"],
            ["Budget", f"EGP {event.budget:,.0f}" if event.budget else "TBD"],
        ]
        table = slide.shapes.add_table(
            len(rows), 2, Inches(0.5), Inches(1.5), Inches(12), Inches(5)
        ).table
        for i, (label, value) in enumerate(rows):
            cell_label = table.cell(i, 0)
            cell_label.text = label
            cell_label.text_frame.paragraphs[0].font.size = Pt(14)
            cell_label.text_frame.paragraphs[0].font.bold = True
            cell_label.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY
            cell_value = table.cell(i, 1)
            cell_value.text = value
            cell_value.text_frame.paragraphs[0].font.size = Pt(14)
            cell_value.text_frame.paragraphs[0].font.color.rgb = self.DARK

    def _add_financial_summary(self, event: Event, line_items: List[SalesLineItem]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY
        header.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.8)
        )
        tf = title.text_frame
        tf.text = "Financial Summary"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        categories = {}
        for item in line_items:
            cat = getattr(item, "category_name", None) or item.description or "General"
            if cat not in categories:
                categories[cat] = {"qty": 0, "total": 0}
            categories[cat]["qty"] += item.quantity or 0
            categories[cat]["total"] += (item.quantity or 0) * (item.unit_price or 0)

        rows = [["Category", "Quantity", "Total"]]
        total = 0
        for cat, data in categories.items():
            rows.append([cat, str(data["qty"]), f"EGP {data['total']:,.0f}"])
            total += data["total"]
        rows.append(["TOTAL", "", f"EGP {total:,.0f}"])

        table = slide.shapes.add_table(
            len(rows), 3, Inches(0.5), Inches(1.5), Inches(12), Inches(5)
        ).table
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = self.WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.PRIMARY

    def _add_timeline_slide(self, event: Event):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY
        header.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.8)
        )
        tf = title.text_frame
        tf.text = "Event Timeline"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        timeline_y = Inches(3.5)
        timeline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), timeline_y, Inches(11), Inches(0.1)
        )
        timeline.fill.solid()
        timeline.fill.fore_color.rgb = self.GRAY
        timeline.line.fill.background()

        milestones = [
            ("Confirmed", event.created_at, self.PRIMARY),
            ("Ops Assigned", getattr(event, "ops_assigned_date", None), self.ACCENT),
            ("Execution", getattr(event, "execution_date", None), self.ACCENT),
            (
                "Completed",
                getattr(event, "completed_date", None),
                RGBColor(0x28, 0xA7, 0x45),
            ),
        ]
        x_pos = Inches(1)
        for label, date, color in milestones:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x_pos - Inches(0.1),
                timeline_y - Inches(0.1),
                Inches(0.2),
                Inches(0.2),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color if date else self.GRAY
            dot.line.fill.background()
            label_box = slide.shapes.add_textbox(
                x_pos - Inches(0.5), timeline_y + Inches(0.3), Inches(1.5), Inches(0.5)
            )
            ltf = label_box.text_frame
            ltf.text = label
            lp = ltf.paragraphs[0]
            lp.font.size = Pt(10)
            lp.font.bold = True
            lp.font.color.rgb = color if date else self.GRAY
            lp.alignment = PP_ALIGN.CENTER
            x_pos += Inches(2.5)

    def _add_vendor_breakdown(self, line_items):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY
        header.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.8)
        )
        tf = title.text_frame
        tf.text = "Vendor Allocation"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

    def _add_checklist_slide(self, event):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY
        header.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.8)
        )
        tf = title.text_frame
        tf.text = "Execution Checklist"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

    def _add_closing_slide(self, event, client):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.PRIMARY
        bg.line.fill.background()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12), Inches(1.5)
        )
        tf = title.text_frame
        tf.text = "Thank You"
        p = tf.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER
        sub = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12), Inches(0.8)
        )
        stf = sub.text_frame
        stf.text = f"IncentiveHouse ERP | Event #{event.id}"
        sp = stf.paragraphs[0]
        sp.font.size = Pt(20)
        sp.font.color.rgb = self.ACCENT
        sp.alignment = PP_ALIGN.CENTER


# ── API Router ──

from fastapi import APIRouter, Depends, BackgroundTasks, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session


router = APIRouter(prefix="/presentations", tags=["Presentations"])


def get_db():
    session = get_sync_session_factory()()
    try:
        yield session
    finally:
        session.close()


class PresentationRequest(BaseModel):
    event_id: int


@router.post("/generate")
def generate_presentation(
    req: PresentationRequest,
    bg: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(require_permission(Permission.REPORT_GENERATE)),
):
    engine = PresentationEngine()
    filepath = engine.generate_event_deck(req.event_id, db)
    return {"status": "generated", "filepath": filepath, "event_id": req.event_id}


import re

_SAFE_FILENAME_RE = re.compile(r"^[a-zA-Z0-9_\-\.]+$")
_MAX_FILENAME_LEN = 255
_PRESENTATIONS_DIR = Path(
    os.getenv("PRESENTATIONS_DIR", str(Path(__file__).parent / "output"))
).resolve()


def _sanitize_filename(raw_name: str) -> str:
    if not raw_name or len(raw_name) > _MAX_FILENAME_LEN:
        raise HTTPException(400, detail="Invalid filename: empty or too long")
    if "/" in raw_name or "\\" in raw_name or ".." in raw_name:
        raise HTTPException(400, detail="Invalid filename: path traversal detected")
    if raw_name.startswith(".") or not _SAFE_FILENAME_RE.match(raw_name):
        raise HTTPException(400, detail="Invalid filename: illegal characters")
    return raw_name


def _safe_file_path(filename: str) -> Path:
    safe_name = _sanitize_filename(filename)
    target = (_PRESENTATIONS_DIR / safe_name).resolve()
    if not str(target).startswith(str(_PRESENTATIONS_DIR)):
        raise HTTPException(400, detail="Invalid filename: path escape detected")
    if not target.exists():
        raise HTTPException(404, detail="File not found")
    return target


@router.get("/download")
def download_presentation(
    path: str,
    current_user: dict = Depends(require_permission(Permission.REPORT_GENERATE)),
):
    safe_path = _safe_file_path(path)
    return FileResponse(
        safe_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=safe_path.name,
    )
